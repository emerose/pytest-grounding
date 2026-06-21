"""Tracked loaders: sha-pinned data + document handles.

``load``/``data`` read a CSV into a DataFrame and record it as provenance; ``doc`` records
a non-table source (PDF/Word/PowerPoint) and returns a :class:`DocRef` whose
:meth:`~DocRef.contains` verifies a verbatim quote against the document's text. All reads
are sha-pinned: the recorded hash is of exactly the bytes that were parsed.

The per-format text readers are offline, deterministic, pure functions of the bytes — no
network, no key, no model. There is **no OCR**: a scanned/image-only document has no text
layer, so rather than let a quote silently "not match", :meth:`DocRef.text` raises
:class:`EmptyExtraction`.
"""
from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path

from ._capture import record
from ._text import match_phrase, preserve_identifier, sha256


class UnsupportedDocFormat(ValueError):
    """Raised by :meth:`DocRef.text` for a suffix no built-in reader handles."""


class EmptyExtraction(RuntimeError):
    """Raised when a document yields no extractable text (e.g. a scanned PDF with no text
    layer). A claim must never quietly pass or fail because its source was unreadable."""


# --------------------------------------------------------------------------- #
# CSV — the tracked data loader
# --------------------------------------------------------------------------- #
def load(path, kind: str = "data"):
    """Read a CSV into a DataFrame, sha-pin it, and record it as provenance.

    The DataFrame carries ``.attrs["source"]`` and ``.attrs["sha256"]``. The sha is of the
    file bytes (exactly what was parsed); the parse goes through ``BytesIO`` so the bypass
    guard never double-counts it. Identifier columns whose values only look numeric (``"01"``,
    ``"08"``) are kept as faithful strings (see :func:`grounding._text.preserve_identifier`).
    Needs the ``[data]`` extra (pandas)."""
    import pandas as pd

    p = Path(path)
    raw = p.read_bytes()
    sha = sha256(raw)
    record(kind, p, sha)
    df = pd.read_csv(io.BytesIO(raw))
    str_df = pd.read_csv(io.BytesIO(raw), dtype=str, keep_default_na=False)
    for col in df.columns:
        df[col] = preserve_identifier(df[col], str_df[col])
    df.attrs["source"] = str(p)
    df.attrs["sha256"] = sha
    return df


data = load  # spelled both ways


# --------------------------------------------------------------------------- #
# Document text readers (the [docs] extra) — offline, deterministic, verbatim
# --------------------------------------------------------------------------- #
def read_pdf_text(path) -> str:
    """All page text of a PDF, newline-joined (pdfplumber). Pure function of the bytes.
    Deliberately not a hosted/Markdown extractor: quote-matching needs raw text that is
    deterministic and verbatim."""
    import pdfplumber

    with pdfplumber.open(str(path)) as pdf:
        return "\n".join((page.extract_text() or "") for page in pdf.pages)


def read_docx_text(path) -> str:
    """All paragraph + table-cell text of a .docx, newline-joined (python-docx)."""
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs]
    for table in d.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def read_pptx_text(path) -> str:
    """All deck prose, newline-joined (python-pptx): title/body text frames, table cells,
    *grouped* shapes, and speaker notes — so a quote in any of them is matchable."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(shape.shapes)
            else:
                yield shape

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in walk(slide.shapes):
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(cell.text for cell in row.cells)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None:
                parts.append(notes.text)
    return "\n".join(parts)


# suffix -> reader. Pure-Python formats only; legacy .doc/.ppt (which would need
# LibreOffice) are intentionally absent and raise UnsupportedDocFormat.
_TEXT_READERS = {
    ".pdf": read_pdf_text,
    ".docx": read_docx_text,
    ".pptx": read_pptx_text,
}

_PRESENTATION_SUFFIXES = {".pptx", ".ppt", ".odp"}


@dataclass
class DocRef:
    """A handle to a non-table source (a PDF/Word report, or a slide deck) recorded as
    evidence. Returned by :func:`doc` so a claim can quote it and keep the citation
    traceable. :meth:`text`/:meth:`contains` extract and match its prose."""

    path: Path
    sha256: str
    _text: str | None = field(default=None, init=False, repr=False, compare=False)

    def __str__(self) -> str:
        return f"{self.path.name}@{self.sha256[:12]}"

    @property
    def is_presentation(self) -> bool:
        """True for slide decks (.pptx/.ppt/.odp) — weaker evidence than a signed report
        (summary text, rounded numbers, scattered across shapes)."""
        return self.path.suffix.lower() in _PRESENTATION_SUFFIXES

    def text(self) -> str:
        """Extract the document's plain text, dispatching on suffix (``.pdf``/``.docx``/
        ``.pptx``; needs the ``[docs]`` extra). Cached on the instance. Raises
        :class:`UnsupportedDocFormat` for any other suffix, and :class:`EmptyExtraction`
        when the document yields no text (scanned/image-only — OCR is not supported)."""
        if self._text is None:
            reader = _TEXT_READERS.get(self.path.suffix.lower())
            if reader is None:
                raise UnsupportedDocFormat(
                    f"doc().text() can't extract {self.path.suffix!r} ({self.path.name}): "
                    f"supported formats are {', '.join(sorted(_TEXT_READERS))} "
                    f"(install the [docs] extra). Legacy .doc/.ppt are not supported.")
            try:
                txt = reader(self.path)
            except ImportError as exc:
                name = getattr(exc, "name", None) or "a reader"
                raise ImportError(
                    f"{name} is required to read {self.path.suffix} — install the [docs] "
                    f"extra: pip install 'grounding[docs]'") from exc
            if not txt.strip():
                raise EmptyExtraction(
                    f"no extractable text in {self.path.name} — a scanned/image-only "
                    f"document? OCR is not supported. A quote can't be verified against "
                    f"an unreadable source.")
            self._text = txt
        return self._text

    def contains(self, phrase: str, *, normalize_ws: bool = True) -> bool:
        """Substring-check ``phrase`` against the extracted :meth:`text`. With
        ``normalize_ws`` (default), fold whitespace/dashes/Markdown on both sides first —
        the robust way to match a verbatim quote an extractor split across lines/cells."""
        return match_phrase(phrase, self.text(), normalize_ws=normalize_ws)


def doc(path, kind: str = "doc"):
    """Record a non-table source (a PDF/Word report, or a slide deck) as a provenance input
    and return a :class:`DocRef`. The quote a claim makes is grounded in the bytes of the
    cited document, sha-pinned like any table. Call :meth:`DocRef.contains` to verify it."""
    p = Path(path)
    sha = sha256(p.read_bytes())
    record(kind, p, sha)
    return DocRef(p, sha)
